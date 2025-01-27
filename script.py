from pptx import Presentation
from pptx.util import Mm
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

# Create presentation
prs = Presentation()

# Set A4 landscape
prs.slide_width = Mm(297)
prs.slide_height = Mm(210)

# Add blank slide
blank_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_layout)

# Set margins and calculate grid
margin = Mm(10)
available_width = prs.slide_width - (2 * margin)
available_height = prs.slide_height - (2 * margin)

# Calculate frame dimensions (10:7 aspect ratio)
frame_width = available_width / 5
frame_height = frame_width * 0.7

# Create 5x4 grid of shapes
for row in range(4):
    for col in range(5):
        left = margin + (col * frame_width)
        top = margin + (row * frame_height)
        
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left,
            top,
            frame_width - Mm(2),
            frame_height - Mm(2)
        )
        
        shape.fill.background()
        shape.line.width = Mm(0.5)

# Create output directory
output_dir = Path.home() / 'Downloads' / 'LogosTemplate'
output_dir.mkdir(parents=True, exist_ok=True)

# Save presentation
output_path = output_dir / 'logo_grid_template.pptx'
prs.save(str(output_path))
print(f"Template saved to: {output_path}")